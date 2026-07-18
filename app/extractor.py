from pathlib import Path
import numpy as np
import pdfplumber
import easyocr
from pdf2image import convert_from_path
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pandas as pd
from PIL import Image

POPPLER_PATH = str(
    Path(__file__).resolve().parent.parent
    / "tools"
    / "poppler-26.02.0"
    / "Library"
    / "bin"
)

reader = None

def extract_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
       return _extract_pdf(file_path)

    elif ext in (".txt", ".md"):
       return _extract_txt(file_path)

    elif ext == ".docx":
       return _extract_docx(file_path)

    elif ext == ".xlsx":
       return _extract_excel(file_path)

    elif ext == ".csv":
       return _extract_csv(file_path)

    elif ext == ".pptx":
       return _extract_ppt(file_path)

    elif ext in (".png", ".jpg", ".jpeg"):
       return _extract_image(file_path)

    else:
       raise ValueError(f"Unsupported file type: {ext}")


def _extract_pdf(file_path: str) -> str:
    text_chunks = []

    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()

            if page_text and page_text.strip():
                text_chunks.append(page_text)

    extracted_text = "\n".join(text_chunks).strip()

    if len(extracted_text) > 50:
        return extracted_text

    return _extract_pdf_with_ocr(file_path)

def _extract_pdf_with_ocr(file_path: str) -> str:
    global reader

    if reader is None:
       reader = easyocr.Reader(["en"], gpu=False)

    pages = convert_from_path(
        file_path,
        dpi=300,
        poppler_path=POPPLER_PATH,
    )

    full_text = []

    for page in pages:
        page_array = np.array(page)
        results = reader.readtext(page_array)

        page_text = " ".join([r[1] for r in results])
        full_text.append(page_text)

    return "\n".join(full_text)


def _extract_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()
    
def _extract_docx(file_path: str) -> str:
    doc = Document(file_path)
    return "\n".join(
        para.text for para in doc.paragraphs if para.text.strip()
    )

def _extract_excel(file_path: str) -> str:
    wb = load_workbook(file_path, data_only=True)

    text = []

    for sheet in wb.worksheets:
        text.append(f"Sheet: {sheet.title}")

        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join(
                str(cell) for cell in row if cell is not None
            )

            if row_text:
                text.append(row_text)

    return "\n".join(text)

def _extract_csv(file_path: str) -> str:
    df = pd.read_csv(file_path, encoding="utf-8", encoding_errors="ignore")

    return df.to_string(index=False)

def _extract_ppt(file_path: str) -> str:
    prs = Presentation(file_path)

    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    text.append(shape.text)

    return "\n".join(text)

def _extract_image(file_path: str) -> str:
    global reader

    if reader is None:
        reader = easyocr.Reader(["en"], gpu=False)

    image = Image.open(file_path).convert("RGB")
    image = np.array(image)

    result = reader.readtext(image)

    return " ".join(r[1] for r in result)